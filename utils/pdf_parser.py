import os
import json
from pypdf import PdfReader
from pptx import Presentation
from utils.gemini_helper import call_gemini, is_api_available

def extract_text_from_pdf(pdf_file_path):
    """
    Extracts text from the uploaded PDF file.
    """
    try:
        reader = PdfReader(pdf_file_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return ""

def extract_text_from_pptx(pptx_file_path):
    """
    Extracts text from the uploaded PPTX file, preserving slide order.
    """
    try:
        prs = Presentation(pptx_file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            if slide_text.strip():
                text += f"--- Slide {i+1} ---\n{slide_text.strip()}\n\n"
        return text.strip()
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return ""

def extract_text_from_file(file_path):
    """
    Router that extracts text based on file extension.
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    return ""

def heuristic_parse_syllabus(pdf_text):
    """
    A smart rule-based fallback parser that extracts units and topics from the actual
    raw PDF text when the Gemini API is not configured or fails.
    """
    lines = [line.strip() for line in pdf_text.split("\n") if line.strip()]
    units = []
    current_unit = None
    
    # Common keywords indicating a unit or chapter
    unit_keywords = ["chapter", "unit", "module", "part", "lecture", "session", "introduction"]
    
    for line in lines:
        line_lower = line.lower()
        is_unit_header = False
        for kw in unit_keywords:
            if line_lower.startswith(kw) or (kw in line_lower and len(line) < 60 and any(char.isdigit() for char in line_lower)):
                is_unit_header = True
                break
                
        # Also check numbered headings like "1. Introduction", "2. Models"
        if not is_unit_header and len(line) < 80:
            parts = line.split()
            if parts and parts[0].endswith(".") and parts[0][:-1].isdigit():
                is_unit_header = True
            elif parts and parts[0].isdigit() and len(parts) > 1 and parts[1].istitle():
                is_unit_header = True
                
        if is_unit_header:
            current_unit = {
                "unit_name": line,
                "topics": []
            }
            units.append(current_unit)
        else:
            if current_unit is not None:
                if len(line) > 5 and len(line) < 120 and len(current_unit["topics"]) < 6:
                    clean_line = line.lstrip("-*+• ").strip()
                    if clean_line and clean_line not in current_unit["topics"]:
                        current_unit["topics"].append(clean_line)
            else:
                # Create initial unit if text starts without a header
                current_unit = {
                    "unit_name": "Unit 1: Overview & Foundation",
                    "topics": []
                }
                units.append(current_unit)
                clean_line = line.lstrip("-*+• ").strip()
                if clean_line:
                    current_unit["topics"].append(clean_line)
                    
    # Clean up empty units or units with no topics
    units = [u for u in units if u["topics"]]
    
    # If we didn't extract anything useful, partition the text as default
    if not units:
        chunk_size = max(1, len(lines) // 4)
        for i in range(0, len(lines), chunk_size):
            chunk_lines = lines[i:i+chunk_size]
            u_name = f"Module {i//chunk_size + 1}: Content Section"
            topics = []
            for cl in chunk_lines:
                clean_cl = cl.lstrip("-*+• ").strip()
                if len(clean_cl) > 10 and len(clean_cl) < 100 and len(topics) < 4:
                    topics.append(clean_cl)
            if topics:
                units.append({
                    "unit_name": u_name,
                    "topics": topics
                })
                
    return {
        "reasoning": f"Rule-based heuristic successfully scanned {len(lines)} lines from the uploaded PDF, dynamically grouping lines based on unit/chapter headers.",
        "units": units[:6]
    }

def parse_syllabus_topics(pdf_text):
    """
    Calls Gemini API (or falls back to mock heuristic) to analyze syllabus text
    and return a structured JSON dict with units and their topics.
    """
    if not pdf_text.strip():
        return {"reasoning": "Empty input text.", "units": []}

    # If API is unavailable, immediately use rule-based actual PDF text parser
    if not is_api_available():
        return heuristic_parse_syllabus(pdf_text)

    # Prepare a prompt for Gemini
    prompt = f"""
    You are an expert curriculum parser. Analyze the following extracted text from a syllabus document and extract its structured units and their corresponding sub-topics. 
    Explain your extraction reasoning in the "reasoning" key.
    
    You must return a raw JSON object. Do not wrap the JSON output in markdown formatting.
    The JSON structure MUST follow this schema:
    {{
      "reasoning": "Explain how units and topics were identified, referencing structural elements.",
      "units": [
        {{
          "unit_name": "Unit Title/Module Name",
          "topics": [
            "Topic description 1",
            "Topic description 2"
          ]
        }}
      ]
    }}

    Syllabus text snippet:
    ---
    {pdf_text[:4000]}
    ---
    """
    
    system_instruction = "You are a syllabus parser that only speaks JSON."
    
    try:
        response_text = call_gemini(prompt, system_instruction=system_instruction, response_mime_type="application/json", raise_exceptions=True)
        clean_text = response_text.replace("```json", "").replace("```", "").strip()
        data = json.loads(clean_text)
        if "units" in data:
            return data
    except Exception as e:
        print(f"Error parsing topics from Gemini response: {e}")
    
    # Fallback to local heuristic extraction instead of static dummy topics
    return heuristic_parse_syllabus(pdf_text)

def merge_syllabi(syllabi_list):
    """
    Merges multiple syllabus structures into a single unified structure.
    Removes duplicate unit names (case-insensitive matches) and duplicate topics.
    """
    merged_units = []
    reasoning_logs = []
    
    for idx, syllabus in enumerate(syllabi_list):
        reasoning_logs.append(f"Merging PDF {idx+1} (found {len(syllabus.get('units', []))} units).")
        for unit in syllabus.get("units", []):
            unit_name = unit.get("unit_name", "").strip()
            unit_topics = unit.get("topics", [])
            if not unit_name:
                continue
                
            # Find if this unit already exists in merged_units
            existing_unit = None
            for mu in merged_units:
                if mu["unit_name"].lower().strip() == unit_name.lower().strip():
                    existing_unit = mu
                    break
            
            if existing_unit:
                # Merge topics, removing duplicates case-insensitively
                for topic in unit_topics:
                    topic_clean = topic.strip()
                    if not any(t.lower().strip() == topic_clean.lower() for t in existing_unit["topics"]):
                        existing_unit["topics"].append(topic_clean)
            else:
                merged_units.append({
                    "unit_name": unit_name,
                    "topics": [t.strip() for t in unit_topics if t.strip()]
                })
                
    merged_units = [mu for mu in merged_units if mu["topics"]]
    
    return {
        "reasoning": " ".join(reasoning_logs) + " Combined units and removed duplicate topics successfully.",
        "units": merged_units
    }



